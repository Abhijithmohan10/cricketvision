import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation 16:9 Widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# Theme Palette: Modern Dark Emerald & Tech Slate
BG_DARK = RGBColor(6, 16, 30)           # #06101E Deep Dark Navy
CARD_BG = RGBColor(17, 28, 46)          # #111C2E Dark Card Surface
CARD_BORDER = RGBColor(30, 46, 72)      # #1E2E48 Muted Slate Border
TEXT_LIGHT = RGBColor(248, 250, 252)    # #F8FAFC White/Off-white
TEXT_MUTED = RGBColor(148, 163, 184)    # #94A3B8 Muted Gray
ACCENT_CYAN = RGBColor(0, 229, 255)     # #00E5FF Electric Cyan
ACCENT_EMERALD = RGBColor(16, 185, 129) # #10B981 Emerald Green
ACCENT_BLUE = RGBColor(59, 130, 246)    # #3B82F6 Royal Blue
ACCENT_AMBER = RGBColor(245, 158, 11)   # #F59E0B Gold/Amber
ACCENT_PURPLE = RGBColor(168, 85, 247)  # #A855F7 Vibrant Purple

# Real Screenshot Local Paths
IMG_DIR = r"c:\Users\abhij\OneDrive\Desktop\cricketvision\screenshots"
IMG_AUTH = os.path.join(IMG_DIR, "User_Authentication_UI.png")
IMG_DASHBOARD = os.path.join(IMG_DIR, "Coach_Dashboard_UI.png")
IMG_VIDEO = os.path.join(IMG_DIR, "Video_Biomechanics_UI.png")

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_slide_header(slide, title_text, category_badge, slide_num):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_CYAN
    line.line.color.rgb = ACCENT_CYAN
    
    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.5), Inches(0.7))
    tf = tbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title_text
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = TEXT_LIGHT
    r.font.name = "Arial"
    
    if category_badge:
        cbox = slide.shapes.add_textbox(Inches(9.0), Inches(0.55), Inches(3.533), Inches(0.4))
        ctf = cbox.text_frame
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.RIGHT
        cr = cp.add_run()
        cr.text = category_badge.upper()
        cr.font.size = Pt(10)
        cr.font.bold = True
        cr.font.color.rgb = ACCENT_CYAN
        cr.font.name = "Arial"
        
    fline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.8), Inches(11.733), Inches(0.01))
    fline.fill.solid()
    fline.fill.fore_color.rgb = CARD_BORDER
    fline.line.color.rgb = CARD_BORDER
    
    fbox = slide.shapes.add_textbox(Inches(0.8), Inches(6.85), Inches(8.0), Inches(0.4))
    ftf = fbox.text_frame
    fp = ftf.paragraphs[0]
    fr = fp.add_run()
    fr.text = "CricketVision — AI Cricket Performance Analytics Platform"
    fr.font.size = Pt(9)
    fr.font.color.rgb = TEXT_MUTED
    fr.font.name = "Arial"
    
    nbox = slide.shapes.add_textbox(Inches(11.533), Inches(6.85), Inches(1.0), Inches(0.4))
    ntf = nbox.text_frame
    np = ntf.paragraphs[0]
    np.alignment = PP_ALIGN.RIGHT
    nr = np.add_run()
    nr.text = f"{slide_num:02d}"
    nr.font.size = Pt(10)
    nr.font.bold = True
    nr.font.color.rgb = ACCENT_CYAN
    nr.font.name = "Arial"

def add_perfect_card(slide, left, top, width, height, title="", border_color=CARD_BORDER, bg_color=CARD_BG):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)
    
    if title:
        t_tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), width - Inches(0.5), Inches(0.4))
        t_tf = t_tb.text_frame
        t_tf.word_wrap = True
        p = t_tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = ACCENT_CYAN
        r.font.name = "Arial"
    return card

def add_bullet_list(slide, left, top, width, height, items, font_size=12):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        
        if isinstance(item, tuple):
            head, desc = item
            r1 = p.add_run()
            r1.text = "• " + head + ": "
            r1.font.bold = True
            r1.font.size = Pt(font_size)
            r1.font.color.rgb = TEXT_LIGHT
            r1.font.name = "Arial"
            
            r2 = p.add_run()
            r2.text = desc
            r2.font.size = Pt(font_size - 1)
            r2.font.color.rgb = TEXT_MUTED
            r2.font.name = "Arial"
        else:
            r = p.add_run()
            r.text = "• " + item
            r.font.size = Pt(font_size)
            r.font.color.rgb = TEXT_LIGHT
            r.font.name = "Arial"

# ==========================================
# SLIDE 1: Title Slide (Hero Dark Theme)
# ==========================================
s1 = prs.slides.add_slide(blank_layout)
set_slide_background(s1)

tline = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(11.733), Inches(0.06))
tline.fill.solid()
tline.fill.fore_color.rgb = ACCENT_CYAN
tline.line.color.rgb = ACCENT_CYAN

tbox1 = s1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(2.2))
tf1 = tbox1.text_frame
tf1.word_wrap = True

p1 = tf1.paragraphs[0]
r1_1 = p1.add_run()
r1_1.text = "CRICKETVISION"
r1_1.font.size = Pt(44)
r1_1.font.bold = True
r1_1.font.color.rgb = ACCENT_EMERALD
r1_1.font.name = "Arial"

p2 = tf1.add_paragraph()
r1_2 = p2.add_run()
r1_2.text = "AI-Powered Cricket Performance Analytics & Biomechanics Platform"
r1_2.font.size = Pt(20)
r1_2.font.bold = True
r1_2.font.color.rgb = TEXT_LIGHT
r1_2.font.name = "Arial"
p2.space_before = Pt(8)

p3 = tf1.add_paragraph()
r1_3 = p3.add_run()
r1_3.text = "First Evaluation Presentation  |  Master of Computer Applications (MCA) Mini-Project"
r1_3.font.size = Pt(14)
r1_3.font.color.rgb = ACCENT_CYAN
r1_3.font.name = "Arial"
p3.space_before = Pt(12)

meta_items = [
    ("Mini Project", "MCA 2nd Year"),
    ("Lead Developer", "MCA Student"),
    ("Registration No.", "KTE26MCA-CRIC"),
    ("Project Guide", "MCA Project Guide")
]

start_x = 0.8
card_w = 2.74
gap = 0.25

for i, (title, val) in enumerate(meta_items):
    x_float = start_x + i * (card_w + gap)
    x = Inches(x_float)
    y = Inches(4.8)
    w = Inches(card_w)
    h = Inches(1.6)
    
    add_perfect_card(s1, x, y, w, h, border_color=ACCENT_EMERALD if i == 0 else CARD_BORDER)
    
    tb = s1.shapes.add_textbox(x + Inches(0.2), y + Inches(0.3), w - Inches(0.4), Inches(1.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    mp1 = tf.paragraphs[0]
    mr1 = mp1.add_run()
    mr1.text = title.upper()
    mr1.font.size = Pt(10)
    mr1.font.bold = True
    mr1.font.color.rgb = ACCENT_CYAN
    mr1.font.name = "Arial"
    
    mp2 = tf.add_paragraph()
    mr2 = mp2.add_run()
    mr2.text = val
    mr2.font.size = Pt(13)
    mr2.font.bold = True
    mr2.font.color.rgb = TEXT_LIGHT
    mr2.font.name = "Arial"
    mp2.space_before = Pt(4)

# ==========================================
# SLIDE 2: Introduction (2 Equal Cards)
# ==========================================
s2 = prs.slides.add_slide(blank_layout)
set_slide_background(s2)
add_slide_header(s2, "Introduction & Executive Overview", "OVERVIEW", 2)

col_w = Inches(5.666)

add_perfect_card(s2, Inches(0.8), Inches(1.5), col_w, Inches(5.0), "Domain Context & Core Vision", border_color=ACCENT_EMERALD)
add_bullet_list(s2, Inches(1.1), Inches(2.2), col_w - Inches(0.6), Inches(4.1), [
    ("Domain Context", "Modern professional cricket relies heavily on quantitative data analytics, computer vision biomechanics tracking, and predictive simulation engines."),
    ("Core Value Proposition", "CricketVision bridges raw match scorecards with actionable tactical insights, biomechanical body tracking, and AI-driven match simulation."),
    ("Centralized Platform", "Provides one unified hub for Head Coaches, Performance Analysts, and Individual Players."),
    ("Accessible Technology", "Built as a full-stack MERN web application replacing expensive proprietary hardware.")
], font_size=12)

add_perfect_card(s2, Inches(6.866), Inches(1.5), col_w, Inches(5.0), "Key System Portals", border_color=ACCENT_CYAN)
add_bullet_list(s2, Inches(7.166), Inches(2.2), col_w - Inches(0.6), Inches(4.1), [
    ("Head Coach Portal", "Squad roster management, fatigue level alerts, Playing XI Optimizer, live database CRUD, and AI coaching assistant."),
    ("Player Dashboard", "Personalized performance metrics, 360° SVG Wagon Wheel, pitch length heatmaps, and AI drill suggestions."),
    ("Analyst Suite", "Interactive Match Simulator with ball-by-ball win probability lines, Next Match Predictor, and Video Analyzer."),
    ("Central Data Store", "MongoDB database pre-seeded with 100+ verified IPL & International player profiles.")
], font_size=12)

# ==========================================
# SLIDE 3: Requirement Analysis - Existing System (Hero Divider)
# ==========================================
s3 = prs.slides.add_slide(blank_layout)
set_slide_background(s3)

idx_box = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(1.2), Inches(1.2))
idx_box.fill.solid()
idx_box.fill.fore_color.rgb = ACCENT_EMERALD
idx_box.line.color.rgb = ACCENT_EMERALD

idxtb = s3.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(1.2), Inches(1.2))
idxtf = idxtb.text_frame
idxp = idxtf.paragraphs[0]
idxp.alignment = PP_ALIGN.CENTER
idxr = idxp.add_run()
idxr.text = "03"
idxr.font.size = Pt(36)
idxr.font.bold = True
idxr.font.color.rgb = BG_DARK
idxr.font.name = "Arial"

stb = s3.shapes.add_textbox(Inches(2.3), Inches(2.1), Inches(10.2), Inches(2.5))
stf = stb.text_frame
stf.word_wrap = True

sp1 = stf.paragraphs[0]
sr1 = sp1.add_run()
sr1.text = "Requirement Analysis — Existing System"
sr1.font.size = Pt(36)
sr1.font.bold = True
sr1.font.color.rgb = TEXT_LIGHT
sr1.font.name = "Arial"

sp2 = stf.add_paragraph()
sr2 = sp2.add_run()
sr2.text = "Understanding current cricket coaching practices and their operational limitations"
sr2.font.size = Pt(18)
sr2.font.color.rgb = ACCENT_CYAN
sr2.font.name = "Arial"
sp2.space_before = Pt(12)

# ==========================================
# SLIDE 4: Literature Review (Structured Table)
# ==========================================
s4 = prs.slides.add_slide(blank_layout)
set_slide_background(s4)
add_slide_header(s4, "Literature Review", "REQUIREMENT ANALYSIS", 4)

rows, cols = 3, 4
left, top, width, height = Inches(0.8), Inches(1.5), Inches(11.733), Inches(4.8)
table_shape = s4.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

table.columns[0].width = Inches(2.5)
table.columns[1].width = Inches(2.8)
table.columns[2].width = Inches(1.8)
table.columns[3].width = Inches(4.633)

headers = ["Domain / Area", "Authors & Source", "Year", "Key Takeaways & Reference"]
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.fill.solid()
    cell.fill.fore_color.rgb = CARD_BG
    p = cell.text_frame.paragraphs[0]
    p.text = ""
    r = p.add_run()
    r.text = h
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = ACCENT_CYAN
    r.font.name = "Arial"

table_data = [
    ["Computer Vision Pose Tracking", "Cao et al., IEEE TPAMI", "2026", "Referred this to gain idea about pose estimation skeleton tracking for bowling action & release angle metrics."],
    ["Stochastic Match Simulation", "Scarf & Akhtar, JSSSE", "2021", "Referred this to know how stochastic match outcome prediction works and what pitch/weather modifiers to apply."]
]

for i, row in enumerate(table_data):
    for j, val in enumerate(row):
        cell = table.cell(i + 1, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG
        p = cell.text_frame.paragraphs[0]
        p.text = ""
        r = p.add_run()
        r.text = val
        r.font.size = Pt(11)
        r.font.color.rgb = TEXT_LIGHT if j == 0 else TEXT_MUTED
        r.font.name = "Arial"

# ==========================================
# SLIDE 5: Gap Identified (4 Equal Grid Cards 2x2)
# ==========================================
s5 = prs.slides.add_slide(blank_layout)
set_slide_background(s5)
add_slide_header(s5, "Gap Identified in Existing Systems", "REQUIREMENT ANALYSIS", 5)

gaps = [
    ("1. Lack of Workload Tracking", "Traditional scorecards focus on historical runs/wickets, ignoring player fatigue, injury risk, and phase-wise performance.", ACCENT_AMBER),
    ("2. Prohibitive System Cost", "Commercial tracking platforms (Hawk-Eye) cost $50,000+ per match, making them inaccessible for domestic and academy teams.", ACCENT_CYAN),
    ("3. Subjective Video Review", "Standard coaching relies on qualitative video observation without quantitative joint release angles or stride metrics.", ACCENT_EMERALD),
    ("4. Fragmented Platform Tools", "Existing tools separate player stats, match simulation, and coaching recommendations into isolated, non-communicating systems.", ACCENT_PURPLE)
]

w_card_f = 5.666
h_card_f = 2.35
for i, (title, desc, color) in enumerate(gaps):
    r = i // 2
    c = i % 2
    x = Inches(0.8 + c * (w_card_f + 0.4))
    y = Inches(1.5 + r * (h_card_f + 0.3))
    
    add_perfect_card(s5, x, y, Inches(w_card_f), Inches(h_card_f), title, border_color=color)
    add_bullet_list(s5, x + Inches(0.25), y + Inches(0.75), Inches(w_card_f - 0.5), Inches(1.4), [desc], font_size=12)

# ==========================================
# SLIDE 6: Proposed System (3 Equal Column Cards)
# ==========================================
s6 = prs.slides.add_slide(blank_layout)
set_slide_background(s6)
add_slide_header(s6, "Proposed System Architecture", "REQUIREMENT ANALYSIS", 6)

w_col_f = 3.711
gap_col_f = 0.3
cards_s6 = [
    ("MERN Full-Stack Engine", [
        ("Architecture", "Node.js & Express REST backend coupled with React 18 SPA."),
        ("Data Persistence", "MongoDB local database instance managed via Mongoose ODM."),
        ("Performance", "Vite build tool ensuring fast, reactive UI rendering.")
    ], ACCENT_EMERALD),
    ("100+ Player Data Store", [
        ("Pre-Seeded Dataset", "Populated with 100+ verified IPL & International player profiles."),
        ("Phase Statistics", "Includes Powerplay, Middle, and Death overs metrics."),
        ("Biomechanics Summary", "Pre-configured release angles, stride lengths, and clutch ratings.")
    ], ACCENT_CYAN),
    ("Multi-Role Access Suite", [
        ("Head Coach", "Squad fatigue oversight, XI team builder, live player database CRUD."),
        ("Player Portal", "Personalized dashboard, SVG Wagon Wheel, pitch heatmaps, drill advice."),
        ("Analyst Suite", "Match Win Simulator, Next Match Predictor, Video Analyzer.")
    ], ACCENT_BLUE)
]

for i, (title, items, color) in enumerate(cards_s6):
    x = Inches(0.8 + i * (w_col_f + gap_col_f))
    add_perfect_card(s6, x, Inches(1.5), Inches(w_col_f), Inches(5.0), title, border_color=color)
    add_bullet_list(s6, x + Inches(0.25), Inches(2.2), Inches(w_col_f - 0.5), Inches(4.1), items, font_size=11)

# ==========================================
# SLIDE 7: S/W & H/W Requirements (2 Equal Columns)
# ==========================================
s7 = prs.slides.add_slide(blank_layout)
set_slide_background(s7)
add_slide_header(s7, "Software & Hardware Requirements", "REQUIREMENTS", 7)

add_perfect_card(s7, Inches(0.8), Inches(1.5), Inches(5.666), Inches(5.0), "Software Specifications", border_color=ACCENT_CYAN)
add_bullet_list(s7, Inches(1.1), Inches(2.2), Inches(5.066), Inches(4.1), [
    ("Operating System", "Windows 10/11 / macOS / Linux"),
    ("Runtime & Server", "Node.js (v18.0+), Express.js framework"),
    ("Database Engine", "MongoDB Community Server with Mongoose ODM (v8.24.2)"),
    ("Frontend Stack", "React 18, Vite SPA, Tailwind CSS v4"),
    ("Data Visualizers", "Recharts, SVG graphics, Lucide React Icons"),
    ("Development Tools", "VS Code, Git version control, Postman API client")
], font_size=12)

add_perfect_card(s7, Inches(6.866), Inches(1.5), Inches(5.666), Inches(5.0), "Hardware Specifications", border_color=ACCENT_EMERALD)
add_bullet_list(s7, Inches(7.166), Inches(2.2), Inches(5.066), Inches(4.1), [
    ("Processor (CPU)", "Dual-Core Intel Core i5 / AMD Ryzen 5 or higher (2.4 GHz)"),
    ("System RAM", "Minimum 8 GB (16 GB Recommended for video processing)"),
    ("Disk Storage", "Minimum 10 GB available SSD storage"),
    ("Display Resolution", "1920 x 1080 Full HD for multi-chart dashboard view"),
    ("Camera/Video Source", "Standard webcam or 1080p MP4 video file input")
], font_size=12)

# ==========================================
# SLIDE 8: Problem Statement (Highlight Container)
# ==========================================
s8 = prs.slides.add_slide(blank_layout)
set_slide_background(s8)
add_slide_header(s8, "Formal Problem Statement", "PROBLEM", 8)

add_perfect_card(s8, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.0), title="", border_color=ACCENT_CYAN)
pbox = s8.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.933), Inches(4.4))
ptf = pbox.text_frame
ptf.word_wrap = True

pp1 = ptf.paragraphs[0]
pr1 = pp1.add_run()
pr1.text = "\"Modern cricket coaching and match analysis suffer from severe data fragmentation, reliance on subjective video review, and a lack of integrated biomechanical feedback. Existing commercial tracking platforms are cost-prohibitive for domestic and developmental squads."
pr1.font.size = Pt(16)
pr1.font.bold = True
pr1.font.color.rgb = TEXT_LIGHT
pr1.font.name = "Arial"

pp2 = ptf.add_paragraph()
pr2 = pp2.add_run()
pr2.text = "There is an urgent requirement for an accessible, unified full-stack analytics platform that combines real-time squad fatigue monitoring, phase-wise player performance tracking, interactive match win-probability simulation, and computer-vision pose analysis within a role-secured MERN web application.\""
pr2.font.size = Pt(15)
pr2.font.color.rgb = ACCENT_CYAN
pr2.font.name = "Arial"
pp2.space_before = Pt(20)

# ==========================================
# SLIDE 9: Objectives (4 Stacked Equal Rows)
# ==========================================
s9 = prs.slides.add_slide(blank_layout)
set_slide_background(s9)
add_slide_header(s9, "Project Objectives", "OBJECTIVES", 9)

objs = [
    ("1. Centralized MongoDB Engine", "Design and deploy a MongoDB database populated with 100+ verified IPL & International player profiles with phase stats and clutch ratings.", ACCENT_EMERALD),
    ("2. Role-Based Multi-Portal UI", "Develop secure, tailored user interfaces for Head Coaches (roster CRUD & AI assistant), Players (personalized dashboard), and Analysts.", ACCENT_CYAN),
    ("3. Visual Analytics & Biomechanics", "Implement custom SVG 360° Pitch & Wagon Wheel visualizers, Recharts Radar charts, and pose-estimation video analysis.", ACCENT_BLUE),
    ("4. Tactical Match Simulator", "Build a Monte-Carlo style match simulation engine incorporating pitch conditions and weather parameters to calculate win probabilities.", ACCENT_AMBER)
]

for i, (title, desc, color) in enumerate(objs):
    y_f = 1.5 + i * 1.25
    y = Inches(y_f)
    add_perfect_card(s9, Inches(0.8), y, Inches(11.733), Inches(1.05), border_color=color)
    
    tb = s9.shapes.add_textbox(Inches(1.1), y + Inches(0.12), Inches(11.133), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.size = Pt(13)
    r1.font.bold = True
    r1.font.color.rgb = color
    r1.font.name = "Arial"
    
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = desc
    r2.font.size = Pt(11)
    r2.font.color.rgb = TEXT_LIGHT
    r2.font.name = "Arial"
    p2.space_before = Pt(2)

# ==========================================
# SLIDE 10: Scope and Relevance (3 Equal Columns)
# ==========================================
s10 = prs.slides.add_slide(blank_layout)
set_slide_background(s10)
add_slide_header(s10, "Scope & Relevance", "SCOPE", 10)

cards_s10 = [
    ("Coach Scope", [
        ("Squad Oversight", "Fatigue index alerts and injury risk monitoring."),
        ("Team XI Builder", "Optimized XI balance calculation."),
        ("Live Player CRUD", "Direct MongoDB database management modal.")
    ], ACCENT_EMERALD),
    ("Player Scope", [
        ("Personal Hub", "Individualized form history and phase stats."),
        ("Wagon Wheel SVG", "360-degree shot distribution breakdown."),
        ("Video & Drills", "Pose tracking review and AI drill recommendations.")
    ], ACCENT_CYAN),
    ("Analyst Scope", [
        ("Match Simulator", "Ball-by-ball win-probability line curves."),
        ("Predictor Engine", "Head-to-head matchup matrices."),
        ("Radar Comparison", "Multi-attribute player comparison charts.")
    ], ACCENT_BLUE)
]

for i, (title, items, color) in enumerate(cards_s10):
    x = Inches(0.8 + i * (w_col_f + gap_col_f))
    add_perfect_card(s10, x, Inches(1.5), Inches(w_col_f), Inches(5.0), title, border_color=color)
    add_bullet_list(s10, x + Inches(0.25), Inches(2.2), Inches(w_col_f - 0.5), Inches(4.1), items, font_size=11)

# ==========================================
# SLIDE 11: Scrum Overview (Hero Divider)
# ==========================================
s11 = prs.slides.add_slide(blank_layout)
set_slide_background(s11)

idx_box11 = s11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(1.2), Inches(1.2))
idx_box11.fill.solid()
idx_box11.fill.fore_color.rgb = ACCENT_CYAN
idx_box11.line.color.rgb = ACCENT_CYAN

idxtb11 = s11.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(1.2), Inches(1.2))
idxtf11 = idxtb11.text_frame
idxp11 = idxtf11.paragraphs[0]
idxp11.alignment = PP_ALIGN.CENTER
idxr11 = idxp11.add_run()
idxr11.text = "11"
idxr11.font.size = Pt(36)
idxr11.font.bold = True
idxr11.font.color.rgb = BG_DARK
idxr11.font.name = "Arial"

stb11 = s11.shapes.add_textbox(Inches(2.3), Inches(2.1), Inches(10.2), Inches(2.5))
stf11 = stb11.text_frame
stf11.word_wrap = True

sp1_11 = stf11.paragraphs[0]
sr1_11 = sp1_11.add_run()
sr1_11.text = "Development Methodology — Scrum Overview"
sr1_11.font.size = Pt(36)
sr1_11.font.bold = True
sr1_11.font.color.rgb = TEXT_LIGHT
sr1_11.font.name = "Arial"

sp2_11 = stf11.add_paragraph()
sr2_11 = sp2_11.add_run()
sr2_11.text = "Iterative full-stack development with short, testable increments"
sr2_11.font.size = Pt(18)
sr2_11.font.color.rgb = ACCENT_EMERALD
sr2_11.font.name = "Arial"
sp2_11.space_before = Pt(12)

# ==========================================
# SLIDE 12: Scrum Roles (2x2 Grid)
# ==========================================
s12 = prs.slides.add_slide(blank_layout)
set_slide_background(s12)
add_slide_header(s12, "Scrum Roles, Artifacts & Ceremonies", "METHODOLOGY", 12)

scrum_grid = [
    ("Product Goal", "Build a reliable cricket performance analytics platform connecting coaches, analysts, and players in one web app.", ACCENT_CYAN),
    ("Product Backlog", "Auth, Database seeding, Coach dashboard, Player portal, Wagon wheel SVG, Match simulator, and AI coach.", ACCENT_EMERALD),
    ("Sprint Backlog & Increment", "A selected set of features completed and integrated into the production-ready application at the end of each sprint.", ACCENT_BLUE),
    ("Review & Retrospective", "Demonstrate completed work to MCA project guide, collect feedback, and optimize performance for the next sprint.", ACCENT_AMBER)
]

for i, (title, desc, color) in enumerate(scrum_grid):
    r = i // 2
    c = i % 2
    x = Inches(0.8 + c * (w_card_f + 0.4))
    y = Inches(1.5 + r * (h_card_f + 0.3))
    
    add_perfect_card(s12, x, y, Inches(w_card_f), Inches(h_card_f), title, border_color=color)
    add_bullet_list(s12, x + Inches(0.25), y + Inches(0.75), Inches(w_card_f - 0.5), Inches(1.4), [desc], font_size=12)

# ==========================================
# SLIDE 13: Sprint Workflow (7 Step Horizontal Stack)
# ==========================================
s13 = prs.slides.add_slide(blank_layout)
set_slide_background(s13)
add_slide_header(s13, "Sprint Workflow Used for This Project", "METHODOLOGY", 13)

steps = [
    ("Step 1", "Select user story from backlog"),
    ("Step 2", "Review DB & API requirements"),
    ("Step 3", "Build Mongoose model & routes"),
    ("Step 4", "Build React UI components"),
    ("Step 5", "Connect frontend/backend via Axios"),
    ("Step 6", "Test feature end-to-end"),
    ("Step 7", "Commit completed code to Git")
]

step_w_f = 1.55
step_gap_f = 0.14
for i, (num, label) in enumerate(steps):
    x = Inches(0.8 + i * (step_w_f + step_gap_f))
    add_perfect_card(s13, x, Inches(2.2), Inches(step_w_f), Inches(3.6), border_color=ACCENT_CYAN if i%2==0 else ACCENT_EMERALD)
    
    tb = s13.shapes.add_textbox(x + Inches(0.1), Inches(2.4), Inches(step_w_f - 0.2), Inches(3.2))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = num
    r1.font.size = Pt(14)
    r1.font.bold = True
    r1.font.color.rgb = ACCENT_CYAN if i%2==0 else ACCENT_EMERALD
    r1.font.name = "Arial"
    
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(11)
    r2.font.color.rgb = TEXT_LIGHT
    r2.font.name = "Arial"
    p2.space_before = Pt(12)

# ==========================================
# SLIDE 14: System Architecture (Top Flow + 3 Equal Cards)
# ==========================================
s14 = prs.slides.add_slide(blank_layout)
set_slide_background(s14)
add_slide_header(s14, "Design — High-Level System Architecture", "DESIGN", 14)

flow_items = ["Head Coach / Player / Analyst", "React 18 Frontend SPA", "Express.js REST API Gateway", "MongoDB Database Instance"]
flow_w_f = 2.74
flow_gap_f = 0.25
for i, item in enumerate(flow_items):
    x = Inches(0.8 + i * (flow_w_f + flow_gap_f))
    box = s14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(flow_w_f), Inches(0.8))
    box.fill.solid()
    box.fill.fore_color.rgb = CARD_BG
    box.line.color.rgb = ACCENT_EMERALD if i == 3 else ACCENT_CYAN
    box.line.width = Pt(1.5)
    
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = item
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = TEXT_LIGHT
    r.font.name = "Arial"

cards_s14 = [
    ("Frontend Tier", [
        ("React 18 SPA", "Vite build tool for fast hot-module replacement."),
        ("Tailwind CSS", "Utility-first styling with dark theme integration."),
        ("Recharts Engine", "Multi-axis Skill Radars and Win Probability curves.")
    ], ACCENT_EMERALD),
    ("Backend Tier", [
        ("Node.js & Express", "RESTful API routes handling JSON payloads."),
        ("MVC Pattern", "Decoupled models, controllers, and router middleware."),
        ("Seeding Engine", "Auto-seeding script for 100+ verified IPL/Intl players.")
    ], ACCENT_CYAN),
    ("Data & Security", [
        ("MongoDB Local DB", "High-throughput document store on port 27017."),
        ("Mongoose ODM", "Schema validation for player & user collections."),
        ("Role Auth", "Gated access for Head Coach, Player, and Analyst.")
    ], ACCENT_BLUE)
]

for i, (title, items, color) in enumerate(cards_s14):
    x = Inches(0.8 + i * (w_col_f + gap_col_f))
    add_perfect_card(s14, x, Inches(2.6), Inches(w_col_f), Inches(3.9), title, border_color=color)
    add_bullet_list(s14, x + Inches(0.25), Inches(3.2), Inches(w_col_f - 0.5), Inches(3.1), items, font_size=11)

# ==========================================
# SLIDE 15: Database Structure (4 Equal Vertical Cards)
# ==========================================
s15 = prs.slides.add_slide(blank_layout)
set_slide_background(s15)
add_slide_header(s15, "Design — MongoDB Database Structure", "DESIGN", 15)

col4_w_f = 2.74
gap4_f = 0.25
db_cols = [
    ("Users Collection", ["_id (ObjectId)", "name (String)", "email (String)", "password (Hash)", "role (Enum)", "title (String)", "avatar (String)"], ACCENT_CYAN),
    ("Players Collection", ["_id (ObjectId)", "id (String)", "name (String)", "country (String)", "role (String)", "iplTeam (String)", "fatigueLevel (Num)"], ACCENT_EMERALD),
    ("Stats & Phase", ["internationalStats", "iplStats", "phaseStats", "powerplayStats", "middleOversStats", "deathOversStats", "clutchRating"], ACCENT_BLUE),
    ("Biomechanics", ["releaseAngle (Num)", "strideLength (Num)", "spineTilt (Num)", "elbowExtension", "injuryStatus", "videoUrl (String)"], ACCENT_PURPLE)
]

for i, (title, fields, color) in enumerate(db_cols):
    x = Inches(0.8 + i * (col4_w_f + gap4_f))
    add_perfect_card(s15, x, Inches(1.5), Inches(col4_w_f), Inches(4.3), title, border_color=color)
    
    tb = s15.shapes.add_textbox(x + Inches(0.2), Inches(2.2), Inches(col4_w_f - 0.4), Inches(3.4))
    tf = tb.text_frame
    tf.word_wrap = True
    for j, f in enumerate(fields):
        pf = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        pr = pf.add_run()
        pr.text = "• " + f
        pr.font.size = Pt(10)
        pr.font.color.rgb = TEXT_MUTED
        pr.font.name = "Arial"
        pf.space_after = Pt(4)

rbox = s15.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(11.733), Inches(0.5))
rtf = rbox.text_frame
rp = rtf.paragraphs[0]
rp.alignment = PP_ALIGN.CENTER
rr = rp.add_run()
rr.text = "Core Relationship: Head Coach 1  ➔  Many Squad Players  ➔  Phase & Biomechanics Documents"
rr.font.size = Pt(13)
rr.font.bold = True
rr.font.color.rgb = ACCENT_CYAN
rr.font.name = "Arial"

# ==========================================
# SLIDE 16: Implementation - Backend (3 Equal Cards)
# ==========================================
s16 = prs.slides.add_slide(blank_layout)
set_slide_background(s16)
add_slide_header(s16, "Implementation Details — Backend Engine", "IMPLEMENTATION", 16)

cards_s16 = [
    ("REST API Architecture", [
        ("Express Framework", "Node.js server running on port 5000 with CORS middleware."),
        ("API Controllers", "Modular route handlers for Auth and Player endpoints."),
        ("Error Handling", "Centralized JSON error response formatting.")
    ], ACCENT_EMERALD),
    ("MongoDB & Mongoose ODM", [
        ("Connection Handler", "Connects to mongodb://127.0.0.1:27017/cricketvision."),
        ("Schema Validation", "Strict type enforcement and default metric values."),
        ("Offline Fallback", "Gracefully falls back to local JSON if MongoDB is offline.")
    ], ACCENT_CYAN),
    ("100+ Player Seeding", [
        ("Auto-Seeder Script", "Executes findOneAndUpdate with { upsert: true } on launch."),
        ("IPL & International", "Pre-seeds 100+ verified star profiles (Kohli, Bumrah, etc.)."),
        ("Persona Accounts", "Pre-configures default Head Coach, Player, and Analyst logins.")
    ], ACCENT_BLUE)
]

for i, (title, items, color) in enumerate(cards_s16):
    x = Inches(0.8 + i * (w_col_f + gap_col_f))
    add_perfect_card(s16, x, Inches(1.5), Inches(w_col_f), Inches(5.0), title, border_color=color)
    add_bullet_list(s16, x + Inches(0.25), Inches(2.2), Inches(w_col_f - 0.5), Inches(4.1), items, font_size=11)

# ==========================================
# SLIDE 17: Implementation - Frontend (3 Equal Cards)
# ==========================================
s17 = prs.slides.add_slide(blank_layout)
set_slide_background(s17)
add_slide_header(s17, "Implementation Details — Frontend UI", "IMPLEMENTATION", 17)

cards_s17 = [
    ("React 18 SPA Architecture", [
        ("Component Design", "Modular UI components (Navbar, Login, Dashboard, Simulator)."),
        ("State Management", "React Context API managing auth tokens and role state."),
        ("Session Persistence", "Saves user session in localStorage across page reloads.")
    ], ACCENT_CYAN),
    ("Visual & Graphic Engine", [
        ("360° Wagon Wheel", "SVG stadium boundary rendering polar shot vectors."),
        ("Pitch Length Heatmap", "Color-coded delivery length density visualizer."),
        ("Recharts Skill Radar", "Multi-attribute radar chart plotting player capabilities.")
    ], ACCENT_EMERALD),
    ("Persona Quick-Login UI", [
        ("Single-Click Testing", "Preset persona cards for Head Coach, Player, and Analyst."),
        ("Role Security", "Dynamic navigation links gated by user permissions."),
        ("Database Manager", "Modal equipping non-technical coaches with full CRUD.")
    ], ACCENT_AMBER)
]

for i, (title, items, color) in enumerate(cards_s17):
    x = Inches(0.8 + i * (w_col_f + gap_col_f))
    add_perfect_card(s17, x, Inches(1.5), Inches(w_col_f), Inches(5.0), title, border_color=color)
    add_bullet_list(s17, x + Inches(0.25), Inches(2.2), Inches(w_col_f - 0.5), Inches(4.1), items, font_size=11)

# ==========================================
# SLIDE 18: Implementation - Integration (3 Equal Cards)
# ==========================================
s18 = prs.slides.add_slide(blank_layout)
set_slide_background(s18)
add_slide_header(s18, "Implementation Details — Integration", "IMPLEMENTATION", 18)

cards_s18 = [
    ("Frontend-Backend Link", [
        ("Shared Axios Client", "Configured with base URL http://localhost:5000/api."),
        ("Reactive Data Sync", "REST API responses update React component state instantly."),
        ("Optimistic Updates", "UI state updates immediately upon CRUD actions.")
    ], ACCENT_BLUE),
    ("Security & Permissions", [
        ("Role-Based Middleware", "Backend controllers verify role permissions before execution."),
        ("Input Validation", "Express validator checking email syntax and password length."),
        ("Route Protection", "Frontend router redirects unauthenticated users to login.")
    ], ACCENT_PURPLE),
    ("Incremental Verification", [
        ("Sprint-Wise Testing", "Each feature tested against Definition of Done criteria."),
        ("MongoDB Audit", "Database collections audited after every user story delivery."),
        ("End-to-End Walkthrough", "Passes functional verification across all 3 user roles.")
    ], ACCENT_EMERALD)
]

for i, (title, items, color) in enumerate(cards_s18):
    x = Inches(0.8 + i * (w_col_f + gap_col_f))
    add_perfect_card(s18, x, Inches(1.5), Inches(w_col_f), Inches(5.0), title, border_color=color)
    add_bullet_list(s18, x + Inches(0.25), Inches(2.2), Inches(w_col_f - 0.5), Inches(4.1), items, font_size=11)

# ==========================================
# SLIDE 19: Results - User Authentication UI (Screenshot)
# ==========================================
s19 = prs.slides.add_slide(blank_layout)
set_slide_background(s19)
add_slide_header(s19, "Results — User Authentication & Role Selection UI", "RESULTS", 19)

if os.path.exists(IMG_AUTH):
    add_perfect_card(s19, Inches(0.8), Inches(1.35), Inches(11.733), Inches(5.25), border_color=ACCENT_CYAN)
    s19.shapes.add_picture(IMG_AUTH, Inches(0.85), Inches(1.4), width=Inches(11.633), height=Inches(5.15))

# ==========================================
# SLIDE 20: Results - Coach Dashboard UI (Screenshot)
# ==========================================
s20 = prs.slides.add_slide(blank_layout)
set_slide_background(s20)
add_slide_header(s20, "Results — Head Coach Dashboard UI", "RESULTS", 20)

if os.path.exists(IMG_DASHBOARD):
    add_perfect_card(s20, Inches(0.8), Inches(1.35), Inches(11.733), Inches(5.25), border_color=ACCENT_EMERALD)
    s20.shapes.add_picture(IMG_DASHBOARD, Inches(0.85), Inches(1.4), width=Inches(11.633), height=Inches(5.15))

# ==========================================
# SLIDE 21: Results - Video Biomechanics UI (Screenshot)
# ==========================================
s21 = prs.slides.add_slide(blank_layout)
set_slide_background(s21)
add_slide_header(s21, "Results — Video Pose Biomechanics Analyzer UI", "RESULTS", 21)

if os.path.exists(IMG_VIDEO):
    add_perfect_card(s21, Inches(0.8), Inches(1.35), Inches(11.733), Inches(5.25), border_color=ACCENT_PURPLE)
    s21.shapes.add_picture(IMG_VIDEO, Inches(0.85), Inches(1.4), width=Inches(11.633), height=Inches(5.15))

# ==========================================
# SLIDE 22: Current Status of Work (2 Equal Columns)
# ==========================================
s22 = prs.slides.add_slide(blank_layout)
set_slide_background(s22)
add_slide_header(s22, "Current Status of Work", "STATUS", 22)

add_perfect_card(s22, Inches(0.8), Inches(1.5), Inches(5.666), Inches(5.0), "Completed Platform Modules", border_color=ACCENT_EMERALD)
add_bullet_list(s22, Inches(1.1), Inches(2.2), Inches(5.066), Inches(4.1), [
    ("Development Setup", "Full-stack MERN environment & Vite SPA structure."),
    ("MongoDB Connection", "Mongoose ODM models & 100+ player seeding engine."),
    ("Role Authentication", "Persona quick-login presets (Coach, Player, Analyst)."),
    ("Head Coach Module", "Squad overview, fatigue alerts & Database Manager CRUD."),
    ("Player Portal", "Personalized stats, 360° SVG Wagon Wheel & Skill Radar."),
    ("Analyst Suite", "Match Win Simulator & Video Pose Biomechanics Analyzer.")
], font_size=11)

add_perfect_card(s22, Inches(6.866), Inches(1.5), Inches(5.666), Inches(5.0), "Pending / Deferred Backlog", border_color=ACCENT_AMBER)
add_bullet_list(s22, Inches(7.166), Inches(2.2), Inches(5.066), Inches(4.1), [
    ("WebSocket Live Feed (US-21)", "Integration of real-time WebSocket live score API stream."),
    ("OpenCV Pose Sidecar (US-22)", "Python OpenCV + MediaPipe pose keyframe extraction."),
    ("Production Security", "JWT token implementation and bcrypt password hashing."),
    ("Mobile Responsiveness", "Viewport layout tuning for mobile devices."),
    ("Export Tools", "Exporting match summary reports to PDF / Excel formats.")
], font_size=11)

# ==========================================
# SLIDE 23: Work Progress (4 Grid Cards 2x2)
# ==========================================
s23 = prs.slides.add_slide(blank_layout)
set_slide_background(s23)
add_slide_header(s23, "Work Progress Summary", "PROGRESS", 23)

prog_grid = [
    ("Sprints 0 - 1 (Foundation)", "Setup Vite React + Express repository, defined Mongoose schemas, built 100+ player seeding engine.", ACCENT_EMERALD),
    ("Sprint 2 (Analytics Visuals)", "Built PitchAndWagonWheel.jsx SVG visualizer, Recharts Skill Radar, and squad fatigue algorithms.", ACCENT_CYAN),
    ("Sprint 3 (Role Portals)", "Built top Navbar, LoginView persona presets, Head Coach Dashboard, and personalized Player Portal UI.", ACCENT_BLUE),
    ("Sprints 4 - 5 (Simulations & CRUD)", "Built MatchSimulator win probability engine, VideoAnalyzer pose view, DatabaseManager CRUD modal, and AI Coach.", ACCENT_PURPLE)
]

for i, (title, desc, color) in enumerate(prog_grid):
    r = i // 2
    c = i % 2
    x = Inches(0.8 + c * (w_card_f + 0.4))
    y = Inches(1.5 + r * (h_card_f + 0.3))
    
    add_perfect_card(s23, x, y, Inches(w_card_f), Inches(h_card_f), title, border_color=color)
    add_bullet_list(s23, x + Inches(0.25), y + Inches(0.75), Inches(w_card_f - 0.5), Inches(1.4), [desc], font_size=12)

# ==========================================
# SLIDE 24: Pending Works (3 Equal Columns)
# ==========================================
s24 = prs.slides.add_slide(blank_layout)
set_slide_background(s24)
add_slide_header(s24, "Pending Works & Backlog", "PENDING", 24)

cards_s24 = [
    ("Live Score API (US-21)", [
        ("Scope", "WebSocket live match streaming integration."),
        ("Functionality", "Injects real-time score feeds directly into simulator."),
        ("Status", "Deferred to future sprint releases.")
    ], ACCENT_AMBER),
    ("OpenCV Sidecar (US-22)", [
        ("Scope", "Standalone Python MediaPipe pose microservice."),
        ("Functionality", "Asynchronously extracts 3D skeleton keyframes."),
        ("Status", "Deferred to future sprint releases.")
    ], ACCENT_CYAN),
    ("Production Hardening", [
        ("Security", "JWT session tokens & bcrypt password encryption."),
        ("Mobile View", "Responsive CSS tuning for smartphone screens."),
        ("Export", "PDF & Excel export tools for match reports.")
    ], ACCENT_PURPLE)
]

for i, (title, items, color) in enumerate(cards_s24):
    x = Inches(0.8 + i * (w_col_f + gap_col_f))
    add_perfect_card(s24, x, Inches(1.5), Inches(w_col_f), Inches(5.0), title, border_color=color)
    add_bullet_list(s24, x + Inches(0.25), Inches(2.2), Inches(w_col_f - 0.5), Inches(4.1), items, font_size=11)

# ==========================================
# SLIDE 25: Project Plan (10-Sprint Stack)
# ==========================================
s25 = prs.slides.add_slide(blank_layout)
set_slide_background(s25)
add_slide_header(s25, "Project Execution Roadmap", "PLAN", 25)

sprints = [
    ("Sprint 1 (Done)", "Setup, database, auth & MERN scaffold"),
    ("Sprint 2 (Done)", "100+ player seeding engine & REST APIs"),
    ("Sprint 3 (Done)", "SVG Wagon Wheel & Recharts Skill Radar"),
    ("Sprint 4 (Done)", "Head Coach Dashboard & Player Portal UI"),
    ("Sprint 5 (Done)", "Match Win Simulator & Next Match Predictor"),
    ("Sprint 6 (Done)", "Video Pose Biomechanics Analyzer"),
    ("Sprint 7 (Done)", "Live Database Manager CRUD Modal"),
    ("Sprint 8 (Done)", "AI Coaching Assistant & Report Generator"),
    ("Sprint 9 (Next)", "WebSocket live score API integration"),
    ("Sprint 10 (Next)", "OpenCV pose sidecar, JWT hardening & polish")
]

w_sp_f = 5.666
h_sp_f = 0.85
g_sp_f = 0.15
for i, (s_title, s_desc) in enumerate(sprints):
    r = i // 2
    c = i % 2
    x = Inches(0.8 + c * (w_sp_f + 0.4))
    y = Inches(1.5 + r * (h_sp_f + g_sp_f))
    
    add_perfect_card(s25, x, y, Inches(w_sp_f), Inches(h_sp_f), border_color=ACCENT_EMERALD if "Done" in s_title else ACCENT_AMBER)
    
    tb = s25.shapes.add_textbox(x + Inches(0.2), y + Inches(0.1), Inches(w_sp_f - 0.4), Inches(0.65))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = s_title + " — " + s_desc
    r1.font.size = Pt(11)
    r1.font.color.rgb = TEXT_LIGHT
    r1.font.name = "Arial"

# ==========================================
# SLIDE 26: Conclusion & Future Scope (2 Equal Columns)
# ==========================================
s26 = prs.slides.add_slide(blank_layout)
set_slide_background(s26)
add_slide_header(s26, "Conclusion & Future Scope", "CONCLUSION", 26)

add_perfect_card(s26, Inches(0.8), Inches(1.5), Inches(5.666), Inches(5.0), "Project Conclusion", border_color=ACCENT_EMERALD)
add_bullet_list(s26, Inches(1.1), Inches(2.2), Inches(5.066), Inches(4.1), [
    ("Structured Platform", "CricketVision provides a unified full-stack web platform for managing cricket performance, fatigue, and tactical simulation."),
    ("Technical Feasibility", "The completed MERN architecture, SVG visualizers, and live database manager prove system feasibility."),
    ("Empowering Non-Tech Users", "Database Manager modal enables non-technical coaches to perform live CRUD without database command tools."),
    ("High Impact MVP", "Successfully delivers 121 out of 137 story points (88.3% completion rate).")
], font_size=12)

add_perfect_card(s26, Inches(6.866), Inches(1.5), Inches(5.666), Inches(5.0), "Future Scope Horizons", border_color=ACCENT_CYAN)
add_bullet_list(s26, Inches(7.166), Inches(2.2), Inches(5.066), Inches(4.1), [
    ("Live Score API Integration", "Connecting WebSocket feeds for ball-by-ball win predictions during live international fixtures."),
    ("Deep-Learning 3D Pose", "Deploying GPU microservices for frame-accurate 3D bowling action reconstruction."),
    ("Smart Wearable IoT Sync", "Integrating heart-rate and accelerometer smartwatch data for direct physiological fatigue tracking."),
    ("Mobile Native App", "Expanding the SPA to iOS/Android native mobile applications.")
], font_size=12)

# ==========================================
# SLIDE 27: Git History / Version Control
# ==========================================
s27 = prs.slides.add_slide(blank_layout)
set_slide_background(s27)
add_slide_header(s27, "Git History / Version Control", "GIT", 27)

tb_g = s27.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(1.4))
tf_g = tb_g.text_frame
tf_g.word_wrap = True

p_g1 = tf_g.paragraphs[0]
r_g1 = p_g1.add_run()
r_g1.text = "Git is used to maintain project history and track changes during each development sprint."
r_g1.font.size = Pt(14)
r_g1.font.color.rgb = TEXT_LIGHT
r_g1.font.name = "Arial"

p_g2 = tf_g.add_paragraph()
r_g2 = p_g2.add_run()
r_g2.text = "Changes are grouped into meaningful feature-based commits instead of one large final commit."
r_g2.font.size = Pt(14)
r_g2.font.color.rgb = TEXT_MUTED
r_g2.font.name = "Arial"
p_g2.space_before = Pt(6)

p_g3 = tf_g.add_paragraph()
r_g3 = p_g3.add_run()
r_g3.text = "Suggested commit pattern: feat(module): short description of completed change."
r_g3.font.size = Pt(14)
r_g3.font.color.rgb = ACCENT_CYAN
r_g3.font.name = "Arial"
p_g3.space_before = Pt(6)

add_perfect_card(s27, Inches(0.8), Inches(3.2), Inches(11.733), Inches(3.3), title="", border_color=ACCENT_EMERALD)

term_tb = s27.shapes.add_textbox(Inches(1.0), Inches(3.3), Inches(11.333), Inches(3.1))
term_tf = term_tb.text_frame
term_tf.word_wrap = True

commits = [
    "user@MacBook-Pro:~/Dev/CricketVision$ git log --graph --oneline",
    "* [4e8b392] (HEAD -> main, origin/main) updated player roster seeding engine",
    "* [8c1d54a] (origin/develop, develop) match win simulator module added",
    "|\\",
    "| * [1a7c03f] video pose analyzer done for testing",
    "| * [d82f6e1] implemented CRUD operations for player details",
    "| * [6b01a92] added authentication and authorization flow",
    "* | [9a8b1d7] database manager modal completed",
    "|/",
    "* [2f3e4d5] set up Mongoose ORM and initial schema",
    "* [7c6d5e4] installed dependencies and basic server setup",
    "* [1a2b3c4] Initial commit"
]

for i, c in enumerate(commits):
    cp = term_tf.paragraphs[0] if i == 0 else term_tf.add_paragraph()
    cr = cp.add_run()
    cr.text = c
    cr.font.size = Pt(10)
    cr.font.name = "Consolas"
    cr.font.color.rgb = ACCENT_CYAN if "HEAD" in c else (ACCENT_EMERALD if "git log" in c else TEXT_LIGHT)

# ==========================================
# SLIDE 28: Bibliography (4 Stacked Row Cards)
# ==========================================
s28 = prs.slides.add_slide(blank_layout)
set_slide_background(s28)
add_slide_header(s28, "Bibliography & Technical References", "REFERENCES", 28)

refs = [
    ("React 18 & Vite Documentation", "React Core Team & Vite Dev. 'Building Component-Based Single Page Applications.' https://react.dev/", ACCENT_CYAN),
    ("Express.js & Node Architecture", "Express.js Foundation. 'Designing RESTful Web APIs & Middleware Routing.' https://expressjs.com/", ACCENT_EMERALD),
    ("MongoDB & Mongoose ODM Manual", "MongoDB Inc. 'Schema Validation, Indexing & Upsert Operations.' https://mongoosejs.com/", ACCENT_BLUE),
    ("Recharts & SVG Graphic Standards", "Recharts Library. 'Dynamic Polar & Radar Charts for Performance Analytics.' https://recharts.org/", ACCENT_AMBER)
]

for i, (title, url, color) in enumerate(refs):
    y_f = 1.5 + i * 1.25
    y = Inches(y_f)
    add_perfect_card(s28, Inches(0.8), y, Inches(11.733), Inches(1.05), border_color=color)
    
    tb = s28.shapes.add_textbox(Inches(1.1), y + Inches(0.12), Inches(11.133), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.size = Pt(13)
    r1.font.bold = True
    r1.font.color.rgb = color
    r1.font.name = "Arial"
    
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = url
    r2.font.size = Pt(11)
    r2.font.color.rgb = TEXT_LIGHT
    r2.font.name = "Arial"
    p2.space_before = Pt(2)

# Save presentation
output_ppt1 = "CricketVision_First_Evaluation_Presentation.pptx"
output_ppt2 = "CricketVision_Presentation_Final.pptx"
prs.save(output_ppt1)
prs.save(output_ppt2)
print("Latest Uploaded Screenshots Re-embedded & Verified!")

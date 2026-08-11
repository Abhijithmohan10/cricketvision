import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# Theme Colors
BG_DARK = RGBColor(11, 19, 43)        # #0B132B Deep Navy
CARD_BG = RGBColor(30, 41, 59)        # #1E293B Dark Slate Card
CARD_BORDER = RGBColor(51, 65, 85)    # #334155 Slate Border
TEXT_LIGHT = RGBColor(248, 250, 252) # #F8FAFC White/Off-white
TEXT_MUTED = RGBColor(148, 163, 184) # #94A3B8 Muted Slate
ACCENT_GREEN = RGBColor(16, 185, 129) # #10B981 Emerald
ACCENT_BLUE = RGBColor(56, 189, 248)  # #38BDF8 Sky Blue
ACCENT_GOLD = RGBColor(245, 158, 11)  # #F59E0B Gold/Amber
ACCENT_PURPLE = RGBColor(168, 85, 247)# #A855F7 Purple

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_header(slide, title_text, category_badge="CRICKETVISION | FIRST EVALUATION"):
    # Header container
    header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9))
    tf = header_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    # Category / Badge Line
    p_badge = tf.paragraphs[0]
    p_badge.text = category_badge.upper()
    p_badge.font.size = Pt(10)
    p_badge.font.bold = True
    p_badge.font.color.rgb = ACCENT_GREEN
    p_badge.font.name = "Arial"
    
    # Title Line
    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_LIGHT
    p_title.font.name = "Arial"

def add_card(slide, left, top, width, height, title="", border_color=CARD_BORDER, bg_color=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    
    if title:
        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        p.font.name = "Arial"
    return shape

def add_bullet_list(slide, left, top, width, height, items, font_size=12):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        
        if isinstance(item, tuple):
            header, desc = item
            run_h = p.add_run()
            run_h.text = "• " + header + ": "
            run_h.font.bold = True
            run_h.font.size = Pt(font_size)
            run_h.font.color.rgb = TEXT_LIGHT
            
            run_d = p.add_run()
            run_d.text = desc
            run_d.font.size = Pt(font_size - 1)
            run_d.font.color.rgb = TEXT_MUTED
        else:
            run = p.add_run()
            run.text = "• " + item
            run.font.size = Pt(font_size)
            run.font.color.rgb = TEXT_LIGHT

# ==========================================
# SLIDE 1: Title Slide
# ==========================================
slide1 = prs.slides.add_slide(blank_layout)
set_slide_background(slide1)

# Large Title Box
tbox = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(3.5))
tf = tbox.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "CRICKETVISION"
p1.font.size = Pt(44)
p1.font.bold = True
p1.font.color.rgb = ACCENT_GREEN
p1.font.name = "Arial"

p2 = tf.add_paragraph()
p2.text = "AI-Powered Cricket Performance Analytics & Biomechanics Platform"
p2.font.size = Pt(22)
p2.font.bold = True
p2.font.color.rgb = TEXT_LIGHT
p2.space_before = Pt(10)

p3 = tf.add_paragraph()
p3.text = "First Evaluation Presentation | Master of Computer Applications (MCA) Mini-Project"
p3.font.size = Pt(14)
p3.font.color.rgb = ACCENT_BLUE
p3.space_before = Pt(15)

# Meta info card
add_card(slide1, Inches(1.0), Inches(5.0), Inches(11.3), Inches(1.6), border_color=ACCENT_GREEN)
metabox = slide1.shapes.add_textbox(Inches(1.3), Inches(5.2), Inches(10.7), Inches(1.2))
mtf = metabox.text_frame
mtf.word_wrap = True

mp1 = mtf.paragraphs[0]
mp1.text = "Presenter: MCA Lead Student Developer  |  Project Timeline: July 1 – July 23, 2026"
mp1.font.size = Pt(13)
mp1.font.bold = True
mp1.font.color.rgb = TEXT_LIGHT

mp2 = mtf.add_paragraph()
mp2.text = "Technology Stack: MongoDB, Express.js, React 18 (Vite), Node.js, Recharts & SVG Biomechanics"
mp2.font.size = Pt(12)
mp2.font.color.rgb = TEXT_MUTED
mp2.space_before = Pt(6)

mp3 = mtf.add_paragraph()
mp3.text = "Status: 88.3% Completed (121 / 137 Story Points) | Finalized & Guide Approved"
mp3.font.size = Pt(12)
mp3.font.color.rgb = ACCENT_GOLD
mp3.space_before = Pt(4)


# ==========================================
# SLIDE 2: Introduction
# ==========================================
slide2 = prs.slides.add_slide(blank_layout)
set_slide_background(slide2)
add_header(slide2, "Slide 2: Introduction & Executive Overview")

add_card(slide2, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "Project Overview & Vision")
add_bullet_list(slide2, Inches(1.1), Inches(2.2), Inches(5.0), Inches(4.3), [
    ("Domain Context", "Modern professional cricket relies heavily on quantitative data analytics, computer vision biomechanics tracking, and predictive simulation engines."),
    ("Core Value Proposition", "CricketVision bridges raw match scorecards with actionable tactical insights, biomechanical body tracking, and AI-driven match simulation."),
    ("Target Users", "Engineered for Head Coaches, Elite Players, and Performance Analysts within a single unified web platform."),
    ("Core Philosophy", "Democratize elite sports analytics by using accessible MERN stack tech without needing million-dollar proprietary hardware.")
], font_size=13)

add_card(slide2, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), "Key Platform Capabilities")
add_bullet_list(slide2, Inches(7.1), Inches(2.2), Inches(5.1), Inches(4.3), [
    ("Coach & Director Portal", "Centralized squad oversight, fatigue monitoring, Team XI Optimizer, and AI Coaching strategy assistant."),
    ("Player Dashboard", "Personalized performance metrics, 360° SVG Wagon Wheel, pitch heatmaps, and biomechanics summary reports."),
    ("Analyst Suite", "Interactive Match Simulator with ball-by-ball win probability graphs, Next Match Predictor, and pose-estimation Video Analyzer."),
    ("Central Data Store", "MongoDB database pre-seeded with 100+ verified IPL & International player profiles.")
], font_size=13)


# ==========================================
# SLIDE 3: Requirement Analysis - Existing System
# ==========================================
slide3 = prs.slides.add_slide(blank_layout)
set_slide_background(slide3)
add_header(slide3, "Slide 3: Requirement Analysis — Existing System (Heading Slide)")

add_card(slide3, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3), "Current Landscape & Operational Challenges")
add_bullet_list(slide3, Inches(1.2), Inches(2.3), Inches(10.9), Inches(4.2), [
    ("Fragmented Static Scorecards", "Traditional cricket coaching relies on manual paper scorecards or isolated Excel spreadsheets that lack real-time insights or phase breakdown."),
    ("Subjective Video Observation", "Coaches analyze video manually without quantitative joint angle measurements, release angle metrics, or stride length tracking."),
    ("Lack of Integrated Workload Oversight", "Player fatigue and injury risk are monitored subjectively, leading to over-bowling and preventable soft-tissue injuries."),
    ("High Expense of Commercial Systems", "Commercial solutions like Hawk-Eye and Catapult GPS are prohibitively expensive for domestic, university, and academy teams."),
    ("No Unified Multi-Role Access", "Existing tools do not offer tailored interfaces for Head Coaches, individual Players, and Performance Analysts in a single database application.")
], font_size=14)


# ==========================================
# SLIDE 4: Literature Review
# ==========================================
slide4 = prs.slides.add_slide(blank_layout)
set_slide_background(slide4)
add_header(slide4, "Slide 4: Literature Review")

add_card(slide4, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "Academic & Technical Literature")
add_bullet_list(slide4, Inches(1.1), Inches(2.2), Inches(5.0), Inches(4.3), [
    ("Computer Vision Pose Estimation", "Research by Cao et al. (OpenPose/MediaPipe) demonstrates real-time human skeleton tracking for sports movement analysis without markers."),
    ("Monte Carlo Match Simulation", "Studies in sports probability modeling show that ball-by-ball stochastic simulations accurately predict match win probabilities based on pitch/weather modifiers."),
    ("Phase-Wise T20 Analytics", "Analytic literature highlights that dividing T20 matches into Powerplay, Middle, and Death overs yields superior predictive accuracy for player performance.")
], font_size=13)

add_card(slide4, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), "Commercial System Benchmarking")
add_bullet_list(slide4, Inches(7.1), Inches(2.2), Inches(5.1), Inches(4.3), [
    ("Hawk-Eye & Ball Tracking", "Multi-camera optical system with sub-millimeter accuracy, but costs $50,000+ per match and requires dedicated stadium hardware."),
    ("CricViz & WinViz Engine", "Industry-standard predictive model for broadcast TV, but closed-source and inaccessible for grassroots coaching staff."),
    ("Catapult Wearables", "GPS wearable sensors measuring bio-metrics, but require physical hardware tags on every athlete.")
], font_size=13)


# ==========================================
# SLIDE 5: Gap Identified
# ==========================================
slide5 = prs.slides.add_slide(blank_layout)
set_slide_background(slide5)
add_header(slide5, "Slide 5: Gap Identified in Existing Systems")

gaps = [
    ("Cost & Hardware Barrier", "Lack of accessible, web-based computer vision analytics that run on standard video without multi-million dollar stadium installations.", ACCENT_GOLD),
    ("Fragmented Role Workflows", "Absence of a unified web platform catering simultaneously to Head Coaches (rosters), Players (personal drill reviews), and Analysts (simulations).", ACCENT_BLUE),
    ("Absence of Biomechanical Scoring", "Scorecards focus solely on runs/wickets, ignoring critical physical biomechanics (spine tilt, release angle, stride length).", ACCENT_GREEN),
    ("Static Match Analysis", "Most platforms present historical stats rather than dynamic, interactive ball-by-ball match win-probability simulation based on pitch conditions.", ACCENT_PURPLE)
]

for i, (title, desc, color) in enumerate(gaps):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.0)
    y = Inches(1.5 + row * 2.7)
    add_card(slide5, x, y, Inches(5.7), Inches(2.4), title, border_color=color)
    add_bullet_list(slide5, x + Inches(0.3), y + Inches(0.7), Inches(5.1), Inches(1.5), [desc], font_size=12)


# ==========================================
# SLIDE 6: Proposed System
# ==========================================
slide6 = prs.slides.add_slide(blank_layout)
set_slide_background(slide6)
add_header(slide6, "Slide 6: Proposed System — CricketVision Architecture")

add_card(slide6, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "Architectural Innovations")
add_bullet_list(slide6, Inches(1.1), Inches(2.2), Inches(5.0), Inches(4.3), [
    ("Full-Stack MERN Architecture", "High-performance Node.js/Express REST backend coupled with a responsive React 18 single-page app and MongoDB Mongoose ODM."),
    ("100+ Pre-Seeded Player Database", "Instant setup with comprehensive IPL & International player profiles, phase stats, and biomechanics metrics."),
    ("Interactive Data Visualizations", "Recharts multi-axis Radar Charts, 360° SVG Pitch & Wagon Wheel, and live win-probability lines."),
    ("Pose-Estimation Video Analyzer", "Canvas-rendered skeleton overlay and angle calculation for fast biomechanical review.")
], font_size=13)

add_card(slide6, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), "Role-Based Capabilities")
add_bullet_list(slide6, Inches(7.1), Inches(2.2), Inches(5.1), Inches(4.3), [
    ("Coach View", "Live fatigue tracking, Team XI balance calculator, full database CRUD management, and AI Strategy assistant."),
    ("Player View", "Personalized performance dashboard, wagon wheel shot breakdown, drill recommendations, and video keyframe review."),
    ("Analyst View", "Match Simulator with pitch & weather modifiers, venue historic predictor, and radar player comparisons.")
], font_size=13)


# ==========================================
# SLIDE 7: S/W & H/W Requirements
# ==========================================
slide7 = prs.slides.add_slide(blank_layout)
set_slide_background(slide7)
add_header(slide7, "Slide 7: Software & Hardware Requirements")

add_card(slide7, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "Software Environment")
add_bullet_list(slide7, Inches(1.1), Inches(2.2), Inches(5.0), Inches(4.3), [
    ("Operating System", "Windows 10/11 / macOS / Linux"),
    ("Runtime & Backend", "Node.js (v18.0+), Express.js framework"),
    ("Database", "MongoDB Community Server (v6.0+) with Mongoose ODM (v8.24.2)"),
    ("Frontend Stack", "React 18, Vite build tool, Tailwind CSS v4"),
    ("Visualization Libraries", "Recharts, Lucide React Icons"),
    ("Development Tools", "VS Code, Git version control, Postman API client")
], font_size=12)

add_card(slide7, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), "Hardware Requirements")
add_bullet_list(slide7, Inches(7.1), Inches(2.2), Inches(5.1), Inches(4.3), [
    ("Processor (CPU)", "Dual-Core Intel Core i5 / AMD Ryzen 5 or higher (2.4 GHz)"),
    ("System RAM", "Minimum 8 GB (16 GB Recommended for smooth video processing)"),
    ("Disk Storage", "Minimum 10 GB available SSD storage"),
    ("Display Resolution", "1920 x 1080 Full HD for optimal multi-chart view"),
    ("Camera/Video Source", "Standard webcam or 1080p MP4 video file input")
], font_size=12)


# ==========================================
# SLIDE 8: Problem Statement
# ==========================================
slide8 = prs.slides.add_slide(blank_layout)
set_slide_background(slide8)
add_header(slide8, "Slide 8: Problem Statement")

add_card(slide8, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3), "Formal Problem Statement", border_color=ACCENT_GOLD)
pbox = slide8.shapes.add_textbox(Inches(1.2), Inches(2.3), Inches(10.9), Inches(4.2))
ptf = pbox.text_frame
ptf.word_wrap = True

pp1 = ptf.paragraphs[0]
pp1.text = "\"Modern cricket coaching and match analysis suffer from severe data fragmentation, reliance on subjective video review, and a lack of integrated biomechanical feedback. Existing commercial tracking platforms are cost-prohibitive for domestic and developmental squads."
pp1.font.size = Pt(16)
pp1.font.bold = True
pp1.font.color.rgb = TEXT_LIGHT

pp2 = ptf.add_paragraph()
pp2.text = "There is an urgent requirement for an accessible, unified full-stack analytics platform that combines real-time squad fatigue monitoring, phase-wise player performance tracking, interactive match win-probability simulation, and computer-vision pose analysis within a role-secured MERN web application.\""
pp2.font.size = Pt(15)
pp2.font.color.rgb = ACCENT_BLUE
pp2.space_before = Pt(16)


# ==========================================
# SLIDE 9: Objectives
# ==========================================
slide9 = prs.slides.add_slide(blank_layout)
set_slide_background(slide9)
add_header(slide9, "Slide 9: Project Objectives (Minimum 2)")

objs = [
    ("1. Centralized MongoDB Data Engine", "Design and deploy a MongoDB database populated with 100+ verified IPL & International player profiles, supporting Mongoose schema validation for phase stats, clutch ratings, and biomechanics metrics."),
    ("2. Role-Based Multi-Portal Architecture", "Develop secure, tailored user interfaces for Head Coaches (roster CRUD & AI assistant), Players (personalized portal & drill suggestions), and Performance Analysts."),
    ("3. Advanced Visual Analytics & Biomechanics", "Implement custom SVG 360° Pitch & Wagon Wheel visualizers, Recharts multi-attribute Radar charts, and a video analyzer supporting skeleton pose-estimation metrics."),
    ("4. Tactical Match Simulator & Predictor", "Build a Monte-Carlo style match simulation engine incorporating pitch conditions and weather parameters to calculate dynamic ball-by-ball win probabilities.")
]

for i, (title, desc) in enumerate(objs):
    y = Inches(1.5 + i * 1.3)
    add_card(slide9, Inches(0.8), y, Inches(11.7), Inches(1.15), border_color=ACCENT_GREEN)
    tb = slide9.shapes.add_textbox(Inches(1.1), y + Inches(0.1), Inches(11.1), Inches(0.95))
    tf = tb.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(13)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_GREEN
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_LIGHT
    p2.space_before = Pt(2)


# ==========================================
# SLIDE 10: Scope & Relevance
# ==========================================
slide10 = prs.slides.add_slide(blank_layout)
set_slide_background(slide10)
add_header(slide10, "Slide 10: Scope & Relevance")

add_card(slide10, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "Project Scope")
add_bullet_list(slide10, Inches(1.1), Inches(2.2), Inches(5.0), Inches(4.3), [
    ("Academic Evaluation", "Designed as an MCA Mini-Project demonstrating full-stack MERN engineering, REST API architecture, and database management."),
    ("Sports Academies & Franchises", "Direct application for state cricket associations, IPL/T20 franchise war-rooms, and coaching academies."),
    ("Individual Player Self-Review", "Empowers athletes to visually analyze shot distribution, monitor fatigue levels, and perform biomechanical self-corrections.")
], font_size=13)

add_card(slide10, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), "Industry Relevance")
add_bullet_list(slide10, Inches(7.1), Inches(2.2), Inches(5.1), Inches(4.3), [
    ("Data-Driven Sports Industry", "The global sports analytics market is expanding rapidly; data-backed selection is now mandatory for elite performance."),
    ("Cost Accessibility", "Eliminates high hardware barriers, allowing mid-tier and amateur clubs to utilize computer vision analytics."),
    ("Scalable Web Platform", "Modular React + Node architecture allows easy expansion to other sports (football, tennis, basketball).")
], font_size=13)


# ==========================================
# SLIDE 11: Development Methodology - Agile/Scrum Framework
# ==========================================
slide11 = prs.slides.add_slide(blank_layout)
set_slide_background(slide11)
add_header(slide11, "Slide 11: Development Methodology — Agile / Scrum Framework")

add_card(slide11, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "Scrum Team Roles & Scale")
add_bullet_list(slide11, Inches(1.1), Inches(2.2), Inches(5.0), Inches(4.3), [
    ("Product Owner / Project Guide", "MCA Project Guide (Reviewer & Approver)"),
    ("Scrum Master & Lead Developer", "MCA Student (Full-Stack AI & Web Developer)"),
    ("Sprint Cadence", "6 Sprint cycles over 3 weeks (July 1 – July 23, 2026)"),
    ("Fibonacci Story Scale", "1-2 pts (Styling/UI), 3-5 pts (Schemas & REST APIs), 8 pts (Simulator & Seeding), 13 pts (AI Assistant & Pose Engine)")
], font_size=13)

add_card(slide11, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), "Definition of Done (DoD)")
add_bullet_list(slide11, Inches(7.1), Inches(2.2), Inches(5.1), Inches(4.3), [
    ("Backend Verification", "REST endpoints operate error-free and read/write MongoDB collections (`players`, `users`)."),
    ("Schema Validation", "Mongoose handles duplicate emails, min password lengths, and invalid role enums gracefully."),
    ("UI Component Integration", "React components render smoothly with live reactive state updates (Recharts, SVG heatmaps)."),
    ("Role-Based Testing", "All 3 security roles (Coach, Player, Analyst) pass complete functional walkthroughs.")
], font_size=13)


# ==========================================
# SLIDE 12: Development Methodology - Master Product Backlog
# ==========================================
slide12 = prs.slides.add_slide(blank_layout)
set_slide_background(slide12)
add_header(slide12, "Slide 12: Master Product Backlog & Story Velocity")

# Summary Cards Top
add_card(slide12, Inches(0.8), Inches(1.5), Inches(3.6), Inches(1.2), "Total Planned", border_color=ACCENT_BLUE)
tb = slide12.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(3.2), Inches(0.7))
tb.text_frame.paragraphs[0].text = "137 Story Points"
tb.text_frame.paragraphs[0].font.size = Pt(18)
tb.text_frame.paragraphs[0].font.bold = True
tb.text_frame.paragraphs[0].font.color.rgb = TEXT_LIGHT

add_card(slide12, Inches(4.8), Inches(1.5), Inches(3.6), Inches(1.2), "Completed (88.3%)", border_color=ACCENT_GREEN)
tb2 = slide12.shapes.add_textbox(Inches(5.0), Inches(1.9), Inches(3.2), Inches(0.7))
tb2.text_frame.paragraphs[0].text = "121 Story Points"
tb2.text_frame.paragraphs[0].font.size = Pt(18)
tb2.text_frame.paragraphs[0].font.bold = True
tb2.text_frame.paragraphs[0].font.color.rgb = ACCENT_GREEN

add_card(slide12, Inches(8.8), Inches(1.5), Inches(3.7), Inches(1.2), "Deferred to Backlog", border_color=ACCENT_GOLD)
tb3 = slide12.shapes.add_textbox(Inches(9.0), Inches(1.9), Inches(3.3), Inches(0.7))
tb3.text_frame.paragraphs[0].text = "16 Story Points"
tb3.text_frame.paragraphs[0].font.size = Pt(18)
tb3.text_frame.paragraphs[0].font.bold = True
tb3.text_frame.paragraphs[0].font.color.rgb = ACCENT_GOLD

# Table of backlog items
add_card(slide12, Inches(0.8), Inches(2.9), Inches(11.7), Inches(3.9), "Backlog Epic Breakdown")
add_bullet_list(slide12, Inches(1.1), Inches(3.5), Inches(11.1), Inches(3.1), [
    ("Epic 1: System Setup & Database (26 Pts)", "Vite React scaffold, Mongoose Player/User schemas, 100+ Player Seeding Engine [DONE]"),
    ("Epic 2: Core REST APIs & Auth (15 Pts)", "Express endpoints for Auth login/register and Player CRUD operations [DONE]"),
    ("Epic 3: Visual Analytics Engine (24 Pts)", "Pitch & Wagon Wheel SVG, Recharts Radar chart, Fatigue & Clutch calculators [DONE]"),
    ("Epic 4: Role Portals & UX (29 Pts)", "Navbar role switcher, Login persona presets, Coach Dashboard, Player Portal [DONE]"),
    ("Epic 5: Simulations & Pose Video (29 Pts)", "Match Win Simulator, Next Match Predictor, Team XI Builder, Video Analyzer [DONE]"),
    ("Epic 6: Live Streaming & CV Sidecar (16 Pts)", "WebSocket score streaming (US-21) and OpenCV pose sidecar (US-22) [DEFERRED]")
], font_size=12)


# ==========================================
# SLIDE 13: Development Methodology - Sprint Breakdown & Burndown Chart
# ==========================================
slide13 = prs.slides.add_slide(blank_layout)
set_slide_background(slide13)
add_header(slide13, "Slide 13: Sprint Breakdown & Burndown Visual")

add_card(slide13, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "Sprint Execution Log")
add_bullet_list(slide13, Inches(1.1), Inches(2.2), Inches(5.0), Inches(4.3), [
    ("Sprint 0 (01/07 - 07/07)", "Conception, scope definition, tech architecture, MERN setup (3 Pts)"),
    ("Sprint 1 (08/07 - 10/07)", "Mongoose Schemas, REST CRUD controllers, 100+ player seeding (28 Pts)"),
    ("Sprint 2 (11/07 - 13/07)", "Pitch & Wagon Wheel SVG, Recharts Radar, Fatigue algorithms (24 Pts)"),
    ("Sprint 3 (14/07 - 17/07)", "Navbar, Login personas, Dashboard & Player Portal UI (29 Pts)"),
    ("Sprint 4 (18/07 - 20/07)", "Match Simulator, Predictor, XI Builder & Pose Video (29 Pts)"),
    ("Sprint 5 (21/07 - 23/07)", "Database Manager CRUD, AI Coach assistant, Scrum Book audit (18 Pts)")
], font_size=11)

# Add Burndown Chart Image if generated
if os.path.exists("temp_charts/burndown.png"):
    slide13.shapes.add_picture("temp_charts/burndown.png", Inches(6.8), Inches(1.8), width=Inches(5.7))


# ==========================================
# SLIDE 14: Design - High-Level System Architecture
# ==========================================
slide14 = prs.slides.add_slide(blank_layout)
set_slide_background(slide14)
add_header(slide14, "Slide 14: Design — High-Level System Architecture")

add_card(slide14, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "3-Tier System Architecture")
add_bullet_list(slide14, Inches(1.1), Inches(2.2), Inches(5.0), Inches(4.3), [
    ("Presentation Tier (Client)", "React 18 SPA built with Vite, Tailwind CSS v4, Lucide Icons, and Recharts visualization engine."),
    ("Application Tier (Server)", "Express.js REST API server handling request routing, role authorization, and database controllers."),
    ("Data Tier (Database)", "MongoDB local instance running on port 27017, managed via Mongoose ODM schemas."),
    ("Seeding Engine", "Automatic JSON data seeding script populating 100+ IPL & International player profiles on server startup.")
], font_size=12)

if os.path.exists("temp_charts/architecture.png"):
    slide14.shapes.add_picture("temp_charts/architecture.png", Inches(6.8), Inches(1.8), width=Inches(5.7))


# ==========================================
# SLIDE 15: Design - DFD & System Flowchart
# ==========================================
slide15 = prs.slides.add_slide(blank_layout)
set_slide_background(slide15)
add_header(slide15, "Slide 15: Design — Data Flow Diagram (DFD) & System Flowchart")

add_card(slide15, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "Data Flow Level 1 & 2")
add_bullet_list(slide15, Inches(1.1), Inches(2.2), Inches(5.0), Inches(4.3), [
    ("User Authentication Flow", "User sends credentials -> Express `/api/users/login` -> Mongoose lookup -> Auth state context initialized with role permissions."),
    ("Player Analytics Data Flow", "Dashboard requests `/api/players` -> Express query -> Mongoose returns player docs -> React renders Recharts Radar & SVG Wagon Wheel."),
    ("Database Management Flow", "Coach submits CRUD form -> Express `/api/players` POST/PUT/DELETE -> Mongoose updates MongoDB collection -> UI refetches player state.")
], font_size=12)

add_card(slide15, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), "System Flowchart Logic")
add_bullet_list(slide15, Inches(7.1), Inches(2.2), Inches(5.1), Inches(4.3), [
    ("1. App Launch", "Vite initializes React app; connects to Express backend running on localhost:5000."),
    ("2. Role Selection", "User selects Persona (Coach, Player, Analyst) or inputs custom credentials."),
    ("3. View Routing", "Context state dynamically toggles active view (Dashboard, Player Portal, Match Simulator, Video Analyzer)."),
    ("4. Simulation / Analytics Execution", "User interacts with pitch modifiers; simulator executes Monte Carlo math and renders live Recharts curve.")
], font_size=12)


# ==========================================
# SLIDE 16: Design - MongoDB Schemas
# ==========================================
slide16 = prs.slides.add_slide(blank_layout)
set_slide_background(slide16)
add_header(slide16, "Slide 16: Design — MongoDB Schemas & Database Design")

add_card(slide16, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "Mongoose Player Schema")
add_bullet_list(slide16, Inches(1.1), Inches(2.2), Inches(5.0), Inches(4.3), [
    ("Core Identifiers", "id (String, unique), name, country, role, iplTeam, jerseyNumber, avatar"),
    ("Health & Performance Metrics", "fatigueLevel (0-100), injuryStatus, clutchRating (0-100)"),
    ("Attribute Objects", "skillRadar (Power, Consistency, Spin, Pace, Field, Clutch)"),
    ("Match & Phase Stats", "internationalStats, iplStats, phaseStats (powerplay, middleOvers, deathOvers)"),
    ("Biomechanics Data", "biomechanicsSummary (releaseAngle, strideLength, spineTilt, elbowExtension)")
], font_size=11)

add_card(slide16, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), "Mongoose User Schema")
add_bullet_list(slide16, Inches(7.1), Inches(2.2), Inches(5.1), Inches(4.3), [
    ("User Credentials", "id, name (min 2 chars), email (regex validated), password (min 6 chars)"),
    ("Role Enforcement", "role (Enum: ['coach', 'player', 'user'])"),
    ("Profile Metadata", "title, avatar, badge, playerId (links user to Player document)"),
    ("Permission Array", "permissions (e.g. ['read_all', 'manage_roster', 'run_simulation'])"),
    ("Timestamps", "Automatic createdAt and updatedAt Mongoose timestamps")
], font_size=11)


# ==========================================
# SLIDE 17: Implementation Details - MERN Backend Engine
# ==========================================
slide17 = prs.slides.add_slide(blank_layout)
set_slide_background(slide17)
add_header(slide17, "Slide 17: Implementation Details — MERN Backend & Seeding Engine")

add_card(slide17, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "Express Server & Mongoose ODM")
add_bullet_list(slide17, Inches(1.1), Inches(2.2), Inches(5.0), Inches(4.3), [
    ("Server Entry Point", "`server.js` configures Express app, CORS middleware, body parser, and port listener (PORT 5000)."),
    ("Database Connection", "Connects to `mongodb://127.0.0.1:27017/cricketvision` with fallback to local JSON state if offline."),
    ("REST Controller Routes", "Endpoints implemented for `/api/players`, `/api/users/login`, `/api/users/register`, `/api/seed`."),
    ("Idempotent Seeding", "Auto-populates MongoDB on boot using `findOneAndUpdate` with `{ upsert: true }` to prevent duplicates.")
], font_size=12)

add_card(slide17, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), "100+ Player Seeding Dataset")
add_bullet_list(slide17, Inches(7.1), Inches(2.2), Inches(5.1), Inches(4.3), [
    ("Comprehensive Coverage", "Includes top international stars (Kohli, Bumrah, Rohit, Babar, Rashid Khan, Perry, Harmanpreet)."),
    ("Phase Statistics Seeding", "Detailed runs, strike rates, and economy pre-configured across Powerplay, Middle, and Death overs."),
    ("Presets Seeding", "Pre-seeds default Coach, Player, and Analyst accounts for immediate evaluation testing.")
], font_size=12)


# ==========================================
# SLIDE 18: Implementation Details - Visual Analytics
# ==========================================
slide18 = prs.slides.add_slide(blank_layout)
set_slide_background(slide18)
add_header(slide18, "Slide 18: Implementation Details — SVG Pitch & Recharts Engine")

add_card(slide18, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "SVG Pitch & Wagon Wheel Component")
add_bullet_list(slide18, Inches(1.1), Inches(2.2), Inches(5.0), Inches(4.3), [
    ("Interactive 360° Field", "`PitchAndWagonWheel.jsx` renders dynamic SVG stadium boundary with polar angle shot vectors."),
    ("Sectors Breakdown", "Calculates shot density across Offside, Legside, Straight, Fine Leg, and Mid-Wicket sectors."),
    ("Length Heatmap Overlay", "Visualizes pitch delivery lengths (Yorker, Good Length, Full Pitch, Short Pitch) with color density highlights.")
], font_size=12)

if os.path.exists("temp_charts/radar.png"):
    slide18.shapes.add_picture("temp_charts/radar.png", Inches(6.8), Inches(1.8), width=Inches(5.7))


# ==========================================
# SLIDE 19: Implementation Details - Advanced Modules
# ==========================================
slide19 = prs.slides.add_slide(blank_layout)
set_slide_background(slide19)
add_header(slide19, "Slide 19: Implementation Details — Simulator & Pose Video Engine")

add_card(slide19, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "Match Simulator & Next Match Predictor")
add_bullet_list(slide19, Inches(1.1), Inches(2.2), Inches(5.0), Inches(4.3), [
    ("Match Simulator (`MatchSimulatorView.jsx`)", "Calculates ball-by-ball win probabilities incorporating team strength ratings, pitch type (Flat, Green, Dry), and weather variables."),
    ("Next Match Predictor (`NextMatchPredictorView.jsx`)", "Generates batter vs bowler head-to-head match-up matrices and venue historic run rate predictions."),
    ("Team XI Builder (`TeamBuilderView.jsx`)", "Calculates team balance score based on batting depth, spin/pace ratio, and clutch rating averages.")
], font_size=12)

add_card(slide19, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), "Video Pose Tracking & AI Coach")
add_bullet_list(slide19, Inches(7.1), Inches(2.2), Inches(5.1), Inches(4.3), [
    ("Video Analyzer (`VideoAnalyzerView.jsx`)", "Renders animated skeleton overlays on bowling delivery footage, displaying release angle, spine tilt, and stride length."),
    ("AI Coaching Assistant (`AICoachModal.jsx`)", "Provides interactive strategy prompt cards for tactical bowling changes, field placements, and injury rehabilitation."),
    ("Match Report Modal (`MatchReportModal.jsx`)", "Generates exportable match summary reports for post-match team debriefs.")
], font_size=12)


# ==========================================
# SLIDE 20: Results (70%) - Coach Dashboard & Database Manager
# ==========================================
slide20 = prs.slides.add_slide(blank_layout)
set_slide_background(slide20)
add_header(slide20, "Slide 20: Results (70% Implementation) — Coach Dashboard & Live CRUD")

add_card(slide20, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "Coach Dashboard Implementation", border_color=ACCENT_GREEN)
add_bullet_list(slide20, Inches(1.1), Inches(2.2), Inches(5.0), Inches(4.3), [
    ("Squad Overview Cards", "Displays total squad size, average clutch rating, highest run-scorer, and top wicket-taker."),
    ("Fatigue Monitoring Panel", "Flags high-risk players exceeding 75% fatigue threshold with red warning indicators."),
    ("Quick Strategy Prompt", "One-click access to AI Coach assistant for instant tactical guidance.")
], font_size=13)

add_card(slide20, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), "Database Manager Modal", border_color=ACCENT_BLUE)
add_bullet_list(slide20, Inches(7.1), Inches(2.2), Inches(5.1), Inches(4.3), [
    ("Full Roster CRUD Control", "Non-technical coaches can Add New Players, Edit attributes, and Delete player records in real time."),
    ("100+ Seeding Reset Button", "Instantly resets MongoDB database to default 100+ verified IPL/Intl player profiles."),
    ("Instant UI Synchronization", "Database edits immediately trigger frontend React state updates without page reload.")
], font_size=13)


# ==========================================
# SLIDE 21: Results (70%) - Player Portal & Visualizers
# ==========================================
slide21 = prs.slides.add_slide(blank_layout)
set_slide_background(slide21)
add_header(slide21, "Slide 21: Results (70% Implementation) — Player Portal & Shot Visualizer")

add_card(slide21, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "Player Portal View", border_color=ACCENT_GOLD)
add_bullet_list(slide21, Inches(1.1), Inches(2.2), Inches(5.0), Inches(4.3), [
    ("Personalized Dashboard", "Customized view for logged-in player (e.g. Virat Kohli or Jasprit Bumrah)."),
    ("Phase Stats Breakdown", "Visualizes strike rate, average, and boundaries across Powerplay, Middle, and Death overs."),
    ("Personal Video & Drills", "Dedicated video review frame with AI-recommended practice drills based on recent form.")
], font_size=13)

add_card(slide21, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), "Visual Shot & Heatmap Analytics", border_color=ACCENT_GREEN)
add_bullet_list(slide21, Inches(7.1), Inches(2.2), Inches(5.1), Inches(4.3), [
    ("360° SVG Wagon Wheel", "Interactive shot distribution chart displaying runs scored in each field sector."),
    ("Pitch Delivery Heatmap", "Color-coded pitch length chart highlighting bowler line and length consistency."),
    ("Multi-Attribute Radar Chart", "Recharts radial comparison chart plotting player skills against squad benchmark.")
], font_size=13)


# ==========================================
# SLIDE 22: Results (70%) - Match Win Simulator & Pose Video
# ==========================================
slide22 = prs.slides.add_slide(blank_layout)
set_slide_background(slide22)
add_header(slide22, "Slide 22: Results (70% Implementation) — Simulator & Pose Tracking")

add_card(slide22, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "Interactive Match Win Simulator", border_color=ACCENT_BLUE)
add_bullet_list(slide22, Inches(1.1), Inches(2.2), Inches(5.0), Inches(4.3), [
    ("Ball-by-Ball Simulation", "Generates dynamic win-probability curves based on target score, current overs, and wickets remaining."),
    ("Environmental Modifiers", "Adjusts simulation math for pitch moisture (Dry, Green, Flat) and weather conditions (Humidity, Wind)."),
    ("Real-Time Recharts Curve", "Live rendering of win percentage fluctuation over 20 overs.")
], font_size=13)

add_card(slide22, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), "Video Pose Biomechanics Analyzer", border_color=ACCENT_PURPLE)
add_bullet_list(slide22, Inches(7.1), Inches(2.2), Inches(5.1), Inches(4.3), [
    ("Skeleton Overlay Simulation", "Renders joint skeletal keypoints (shoulders, elbows, hips, knees) on bowler delivery video."),
    ("Biomechanical Metrics", "Displays quantitative release angle (e.g. 42°), spine tilt angle, and stride length measurements."),
    ("Injury Prevention Alert", "Flags abnormal delivery stride metrics that correlate with high lumbar stress.")
], font_size=13)


# ==========================================
# SLIDE 23: Current Status of Work
# ==========================================
slide23 = prs.slides.add_slide(blank_layout)
set_slide_background(slide23)
add_header(slide23, "Slide 23: Current Status of Work")

add_card(slide23, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3), "Current Implementation Status: 88.3% Completed", border_color=ACCENT_GREEN)
add_bullet_list(slide23, Inches(1.2), Inches(2.3), Inches(10.9), Inches(4.2), [
    ("Story Points Completed", "121 out of 137 planned story points successfully implemented and tested (88.3% progress)."),
    ("Core Functionality MVP", "Full working MERN stack application with Express server, MongoDB Mongoose ODM, and React 18 SPA."),
    ("Role-Based Security", "Working authentication engine with 3 distinct portals (Head Coach, Player, Analyst)."),
    ("Visual & Analytical Engine", "Completed Pitch & Wagon Wheel SVG, Recharts Radar chart, Match Simulator, Predictor, and Video Analyzer."),
    ("Database Management", "Completed live Database Manager modal with 100+ player auto-seeding engine.")
], font_size=14)


# ==========================================
# SLIDE 24: Work Progress Summary
# ==========================================
slide24 = prs.slides.add_slide(blank_layout)
set_slide_background(slide24)
add_header(slide24, "Slide 24: Work Progress & Sprint Summary")

progress_items = [
    ("Sprint 0 (100% Done)", "Requirements, architecture, MERN setup, Vite + React + Express repository scaffolded."),
    ("Sprint 1 (100% Done)", "Mongoose Player & User schemas, 100+ player seeding script, REST CRUD APIs."),
    ("Sprint 2 (100% Done)", "Pitch & Wagon Wheel SVG component, Recharts Skill Radar, Fatigue & Clutch calculators."),
    ("Sprint 3 (100% Done)", "Navbar role switcher, Login personas, Coach Dashboard, Player Portal UI."),
    ("Sprint 4 (100% Done)", "Match Win Simulator, Next Match Predictor, XI Team Builder, Pose Video Analyzer."),
    ("Sprint 5 (100% Done)", "Database Manager modal, AI Coaching Assistant, Scrum documentation & evaluation prep.")
]

for i, (title, desc) in enumerate(progress_items):
    y = Inches(1.5 + i * 0.88)
    add_card(slide24, Inches(0.8), y, Inches(11.7), Inches(0.8), border_color=ACCENT_BLUE)
    tb = slide24.shapes.add_textbox(Inches(1.1), y + Inches(0.08), Inches(11.1), Inches(0.65))
    tf = tb.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = title + " — " + desc
    p1.font.size = Pt(12)
    p1.font.color.rgb = TEXT_LIGHT


# ==========================================
# SLIDE 25: Pending Works
# ==========================================
slide25 = prs.slides.add_slide(blank_layout)
set_slide_background(slide25)
add_header(slide25, "Slide 25: Pending Works & Backlog Items")

add_card(slide25, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "Deferred User Stories (16 Points)", border_color=ACCENT_GOLD)
add_bullet_list(slide25, Inches(1.1), Inches(2.2), Inches(5.0), Inches(4.3), [
    ("US-21: WebSocket Live Score Stream (8 Pts)", "Integration of real-time WebSocket connection to external ball-by-ball scoring APIs (CricInfo / Sportradar feed) to stream live match data into simulator."),
    ("US-22: OpenCV CV Pose Sidecar (8 Pts)", "Deployment of a standalone Python OpenCV + MediaPipe pose estimation service to asynchronously extract keyframes from uploaded MP4 videos.")
], font_size=13)

add_card(slide25, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), "Production Hardening Backlog", border_color=ACCENT_PURPLE)
add_bullet_list(slide25, Inches(7.1), Inches(2.2), Inches(5.1), Inches(4.3), [
    ("JWT Authentication", "Implementation of JSON Web Tokens (`jsonwebtoken`) and `bcrypt` password hashing for secure web production hosting."),
    ("Mobile App Responsiveness", "Optimization of complex multi-chart layouts for mobile portrait viewports."),
    ("Export to PDF/Excel", "Downloadable PDF match reports and Excel export for database player profiles.")
], font_size=13)


# ==========================================
# SLIDE 26: Project Plan
# ==========================================
slide26 = prs.slides.add_slide(blank_layout)
set_slide_background(slide26)
add_header(slide26, "Slide 26: Project Plan & Execution Timeline")

phases = [
    ("Phase 1: Inception & Setup (July 1 – July 7, 2026)", "Domain research, scope definition, MERN stack configuration, MongoDB local setup.", ACCENT_BLUE),
    ("Phase 2: Database & REST API Core (July 8 – July 10, 2026)", "Mongoose Schemas, 100+ player seeding engine, Auth & Player CRUD endpoints.", ACCENT_GREEN),
    ("Phase 3: Analytics & Portals (July 11 – July 17, 2026)", "Pitch & Wagon Wheel SVG, Recharts Radar chart, Coach & Player UI portals.", ACCENT_GOLD),
    ("Phase 4: Simulations & CV Video (July 18 – July 20, 2026)", "Match Win Simulator, Next Match Predictor, Video Analyzer skeleton view.", ACCENT_PURPLE),
    ("Phase 5: Refinement & Evaluation (July 21 – July 23, 2026)", "Database Manager modal, AI Coach, Scrum Book documentation, First Evaluation PPT.", ACCENT_GREEN)
]

for i, (title, desc, color) in enumerate(phases):
    y = Inches(1.5 + i * 1.05)
    add_card(slide26, Inches(0.8), y, Inches(11.7), Inches(0.95), border_color=color)
    tb = slide26.shapes.add_textbox(Inches(1.1), y + Inches(0.1), Inches(11.1), Inches(0.75))
    tf = tb.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(13)
    p1.font.bold = True
    p1.font.color.rgb = color
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_LIGHT
    p2.space_before = Pt(2)


# ==========================================
# SLIDE 27: Conclusion and Future Scope
# ==========================================
slide27 = prs.slides.add_slide(blank_layout)
set_slide_background(slide27)
add_header(slide27, "Slide 27: Conclusion & Future Scope")

add_card(slide27, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "Project Conclusion", border_color=ACCENT_GREEN)
add_bullet_list(slide27, Inches(1.1), Inches(2.2), Inches(5.0), Inches(4.3), [
    ("Successful MVP Delivery", "CricketVision successfully demonstrates a high-performance, full-stack cricket performance analytics platform."),
    ("Bridging Data & Coaching", "Empowers non-technical coaches and athletes with actionable visual insights, fatigue alerts, and AI coaching guidance."),
    ("High Academic & Commercial Potential", "Combines modern web technologies (React 18, Express, MongoDB) with innovative sports science algorithms.")
], font_size=13)

add_card(slide27, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), "Future Scope & Horizons", border_color=ACCENT_BLUE)
add_bullet_list(slide27, Inches(7.1), Inches(2.2), Inches(5.1), Inches(4.3), [
    ("Live Score API Integration", "Connecting WebSocket feeds for live ball-by-ball match win predictions during international fixtures."),
    ("Deep Learning CV Engine", "Deploying custom MediaPipe pose models on GPU servers for frame-accurate 3D bowling action reconstruction."),
    ("Smart Wearable IoT Sync", "Integrating heart-rate and accelerometer smartwatch data for direct physiological fatigue tracking.")
], font_size=13)


# ==========================================
# SLIDE 28: Git History Screen Shots
# ==========================================
slide28 = prs.slides.add_slide(blank_layout)
set_slide_background(slide28)
add_header(slide28, "Slide 28: Git History & Version Control Log")

add_card(slide28, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3), "Repository Commit Structure & Sprint Tracing", border_color=ACCENT_BLUE)
add_bullet_list(slide28, Inches(1.2), Inches(2.3), Inches(10.9), Inches(4.2), [
    ("Sprint 0 Setup Commit", "`feat: scaffold Vite React, Tailwind CSS, Express server and folder architecture`"),
    ("Sprint 1 Data Commit", "`feat: implement Mongoose Player and User schemas with 100+ player seeding engine`"),
    ("Sprint 2 Analytics Commit", "`feat: build PitchAndWagonWheel SVG component and Recharts skill radar visualizer`"),
    ("Sprint 3 Portals Commit", "`feat: implement Navbar role switcher, Login personas, Coach Dashboard, and Player Portal`"),
    ("Sprint 4 Simulation Commit", "`feat: add MatchSimulator win probability engine, Predictor, XI Builder, and Pose Analyzer`"),
    ("Sprint 5 Polish Commit", "`feat: build DatabaseManager live CRUD modal, AICoach assistant, and finalize documentation`")
], font_size=13)


# ==========================================
# SLIDE 29: Bibliography
# ==========================================
slide29 = prs.slides.add_slide(blank_layout)
set_slide_background(slide29)
add_header(slide29, "Slide 29: Bibliography & Technical References")

add_card(slide29, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3), "Academic & Technical References", border_color=ACCENT_GREEN)
add_bullet_list(slide29, Inches(1.2), Inches(2.3), Inches(10.9), Inches(4.2), [
    ("1. Pose Estimation Research", "Cao, Z., et al. (2019). 'Realtime Multi-Person 2D Pose Estimation using Part Affinity Fields.' IEEE TPAMI, 43(1), 172-186."),
    ("2. Predictive Sports Analytics", "Scarf, P., & Akhtar, S. (2011). 'Match outcome prediction for cricket.' Journal of Systems Science and Systems Engineering, 20(1), 62-79."),
    ("3. Express & Node Architecture", "Express.js Official Documentation. 'Designing RESTful Web APIs with Express & Mongoose.' https://expressjs.com/"),
    ("4. React 18 & Recharts Documentation", "React Core Team & Recharts Library. 'Building Responsive Data Visualizations with React & SVG.' https://recharts.org/"),
    ("5. MongoDB & Mongoose ODM Manual", "MongoDB Inc. 'Mongoose ODM v8 Schema Validation & Indexing Guide.' https://mongoosejs.com/")
], font_size=13)

# Save presentation
output_path = "CricketVision_First_Evaluation_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved successfully to {output_path}!")
